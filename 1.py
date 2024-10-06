import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from googletrans import Translator, LANGUAGES
import threading
import speech_recognition as sr
from gtts import gTTS
from playsound import playsound
import pytesseract
from PIL import Image
import os
import fitz
import docx
import moviepy.editor as mp
from pptx import Presentation

class LiveMultimodalTranslator:
    def __init__(self):
        self.translator = Translator()
        self.setup_gui()
        self.populate_language_combos()
        self.set_tesseract_path()
        self.root.mainloop()

    def setup_gui(self):
        self.colors = {
            "bg": "#f0f0f0",
            "button": "#4a4e69",
            "button_hover": "#9a8c98",
            "input_bg": "#dcdcdc",
            "output_bg": "#dcdcdc",
            "text": "#333333",
            "button_text": "#ffffff"
        }
        self.root = tk.Tk()
        self.root.title("Live Multimodal Language Translation System")
        self.root.geometry("1200x800")
        self.root.resizable(False, False)
        self.root.config(bg=self.colors["bg"])
        self.title_label = tk.Label(self.root, text="LIVE MULTIMODAL LANGUAGE TRANSLATION SYSTEM", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 20, "bold"))
        self.title_label.pack(pady=20)
        self.main_frame = tk.Frame(self.root, bg=self.colors["bg"])
        self.main_frame.pack(pady=20)
        self.setup_input_frame()
        self.setup_output_frame()
        self.setup_controls_frame()

    def setup_input_frame(self):
        self.input_frame = tk.LabelFrame(self.main_frame, text="INPUT", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12, "bold"), padx=20, pady=20)
        self.input_frame.grid(row=0, column=0, padx=20, pady=10)
        self.input_label = tk.Label(self.input_frame, text="", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12))
        self.input_label.grid(row=0, column=0, sticky="w")
        self.original_text = tk.Text(self.input_frame, height=10, width=40, bg=self.colors["input_bg"], font=("Times New Roman", 15))
        self.original_text.grid(row=1, column=0, pady=10)
        self.source_label = tk.Label(self.input_frame, text="Source language:", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12))
        self.source_label.grid(row=2, column=0, sticky="w")
        self.original_combo = ttk.Combobox(self.input_frame, width=50)
        self.original_combo.grid(row=3, column=0, pady=10)
        self.voice_input_button = tk.Button(self.input_frame, text="VOICE INPUT", font=("Helvetica", 12), command=self.start_voice_input, bg=self.colors["button"], fg=self.colors["button_text"])
        self.voice_input_button.grid(row=4, column=0, pady=10, sticky="ew")
        self.image_input_button = tk.Button(self.input_frame, text="UPLOAD PHOTO", font=("Helvetica", 12), command=self.start_load_image, bg=self.colors["button"], fg=self.colors["button_text"])
        self.image_input_button.grid(row=5, column=0, pady=10, sticky="ew")
        self.file_input_button = tk.Button(self.input_frame, text="UPLOAD DOCUMENT", font=("Helvetica", 12), command=self.start_load_file, bg=self.colors["button"], fg=self.colors["button_text"])
        self.file_input_button.grid(row=6, column=0, pady=10, sticky="ew")

    def setup_output_frame(self):
        self.output_frame = tk.LabelFrame(self.main_frame, text="OUTPUT", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12, "bold"), padx=20, pady=20)
        self.output_frame.grid(row=0, column=1, padx=20, pady=10)
        self.output_label = tk.Label(self.output_frame, text="", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12))
        self.output_label.grid(row=0, column=0, sticky="w")
        self.translated_text = tk.Text(self.output_frame, height=10, width=40, bg=self.colors["output_bg"], font=("Times New Roman", 15))
        self.translated_text.grid(row=1, column=0, pady=10)
        self.voice_output_button = tk.Button(self.output_frame, text="VOICE OUTPUT", font=("Helvetica", 12), command=self.start_voice_output, bg=self.colors["button"], fg=self.colors["button_text"])
        self.voice_output_button.grid(row=2, column=0, pady=10, sticky="ew")
        self.image_label = tk.Label(self.output_frame, text="Image:", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12))
        self.image_label.grid(row=3, column=0, sticky="w")
        self.target_label = tk.Label(self.output_frame, text="Target language:", bg=self.colors["bg"], fg=self.colors["text"], font=("Helvetica", 12))
        self.target_label.grid(row=4, column=0, sticky="w")
        self.translated_combo = ttk.Combobox(self.output_frame, width=50)
        self.translated_combo.grid(row=5, column=0, pady=10)

    def setup_controls_frame(self):
        self.controls_frame = tk.Frame(self.main_frame, bg=self.colors["bg"])
        self.controls_frame.grid(row=0, column=2, padx=20, pady=10)
        self.translate_button = tk.Button(self.controls_frame, text="TRANSLATE", font=("Helvetica", 12), command=self.start_translation, bg="#58a4b0", fg=self.colors["button_text"])
        self.translate_button.grid(row=0, column=0, pady=10, sticky="ew")
        self.copy_button = tk.Button(self.controls_frame, text="COPY TEXT", font=("Helvetica", 12), command=self.copy_to_clipboard, bg="#58a4b0", fg=self.colors["button_text"])
        self.copy_button.grid(row=1, column=0, pady=10, sticky="ew")
        self.clear_button = tk.Button(self.controls_frame, text="CLEAR", font=("Helvetica", 12), command=self.clear_texts, bg="#d92027", fg=self.colors["button_text"])
        self.clear_button.grid(row=2, column=0, pady=10, sticky="ew")
        self.progress = ttk.Progressbar(self.controls_frame, orient="horizontal", length=200, mode="determinate")
        self.progress.grid(row=3, column=0, pady=10, sticky="ew")

    def populate_language_combos(self):
        language_list = list(LANGUAGES.values())
        self.original_combo['values'] = language_list
        self.translated_combo['values'] = language_list
        self.original_combo.set("english")
        self.translated_combo.set("spanish")

    def set_tesseract_path(self):
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        os.environ["TESSDATA_PREFIX"] = r"C:\Program Files\Tesseract-OCR\tessdata"

    def start_translation(self):
        threading.Thread(target=self.translate_it).start()

    def start_voice_input(self):
        threading.Thread(target=self.voice_input).start()

    def start_voice_output(self):
        threading.Thread(target=self.voice_output).start()

    def start_load_image(self):
        threading.Thread(target=self.load_image).start()

    def start_load_file(self):
        threading.Thread(target=self.load_file).start()

    def translate_it(self):
        self.progress.start(10)
        original_text = self.original_text.get("1.0", tk.END)
        src_lang = self.original_combo.get()
        dest_lang = self.translated_combo.get()

        if not original_text.strip():
            self.show_message("Please provide text for translation.")
            self.progress.stop()
            return

        try:
            translation = self.translator.translate(original_text, src=src_lang, dest=dest_lang)
            translated_text = translation.text
        except Exception as e:
            self.show_message(f"Error during translation: {e}")
            self.progress.stop()
            return

        self.translated_text.delete("1.0", tk.END)
        self.translated_text.insert("1.0", translated_text)
        self.progress.stop()

    def voice_input(self):
        self.progress.start(10)
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            self.show_message("Listening...")
            audio = recognizer.listen(source)

        try:
            text = recognizer.recognize_google(audio)
            self.original_text.insert(tk.END, text)
        except sr.UnknownValueError:
            self.show_message("Sorry, I could not understand the audio.")
        except sr.RequestError as e:
            self.show_message(f"Could not request results; {e}")

        self.progress.stop()

    def voice_output(self):
        self.progress.start(10)
        text = self.translated_text.get("1.0", tk.END)
        tts = gTTS(text)
        temp_file = "temp_output.mp3"
        tts.save(temp_file)
        playsound(temp_file)
        os.remove(temp_file)
        self.progress.stop()

    def load_image(self):
        file_path = filedialog.askopenfilename()
        if not file_path:
            return

        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            self.original_text.insert(tk.END, text)
        except Exception as e:
            self.show_message(f"Error loading image: {e}")

    def load_file(self):
        file_path = filedialog.askopenfilename()
        if not file_path:
            return

        extension = file_path.split(".")[-1]
        try:
            if extension == "pdf":
                self.read_pdf(file_path)
            elif extension == "docx":
                self.read_docx(file_path)
            elif extension in ["pptx", "ppt"]:
                self.read_pptx(file_path)
            elif extension in ["mp4", "mkv", "avi"]:
                self.read_video(file_path)
            else:
                self.show_message("Unsupported file format.")
        except Exception as e:
            self.show_message(f"Error loading file: {e}")

    def read_pdf(self, path):
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text()
        self.original_text.insert(tk.END, text)

    def read_docx(self, path):
        doc = docx.Document(path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        self.original_text.insert(tk.END, text)

    def read_pptx(self, path):
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"
        self.original_text.insert(tk.END, text)

    def read_video(self, path):
        video = mp.VideoFileClip(path)
        audio = video.audio
        audio.write_audiofile("temp_audio.wav")
        recognizer = sr.Recognizer()
        with sr.AudioFile("temp_audio.wav") as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data)
        os.remove("temp_audio.wav")
        self.original_text.insert(tk.END, text)

    def copy_to_clipboard(self):
        self.root.clipboard_clear()
        self.root.clipboard_append(self.translated_text.get("1.0", tk.END).strip())
        self.root.update()

    def clear_texts(self):
        self.original_text.delete("1.0", tk.END)
        self.translated_text.delete("1.0", tk.END)

    def show_message(self, message):
        messagebox.showinfo("Information", message)

if __name__ == "__main__":
    app = LiveMultimodalTranslator()
